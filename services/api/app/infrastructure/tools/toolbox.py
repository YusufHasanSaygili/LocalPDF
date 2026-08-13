from __future__ import annotations

import difflib
import html
import io
import os
import re
import sys
import zipfile
from collections import Counter
from collections.abc import Iterable
from pathlib import Path
from typing import Any, cast

import ctranslate2  # type: ignore[import-untyped]
import numpy as np
import pikepdf
import pymupdf  # type: ignore[import-untyped]
import sentencepiece as sentencepiece
from bs4 import BeautifulSoup
from docx import Document as WordDocument
from openpyxl import Workbook, load_workbook  # type: ignore[import-untyped]
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches
from rapidocr import RapidOCR
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.domain.errors import LocalPDFError
from app.domain.validation import parse_page_ranges
from app.infrastructure.tools.pdf import validate_output


def _pages(expression: str, count: int) -> list[int]:
    return [page for group in parse_page_ranges(expression, count) for page in group]


def remove_pages(source: Path, output: Path, expression: str) -> int:
    with pikepdf.open(source) as pdf:
        removed = set(_pages(expression, len(pdf.pages)))
        kept = [page for number, page in enumerate(pdf.pages, start=1) if number not in removed]
        if not kept:
            raise LocalPDFError("INPUT_INVALID", "At least one page must remain.")
        result = pikepdf.Pdf.new()
        result.pages.extend(kept)
        result.save(output)
    return validate_output(output, len(kept))


def extract_pages(source: Path, output: Path, expression: str) -> int:
    with pikepdf.open(source) as pdf:
        selected = _pages(expression, len(pdf.pages))
        result = pikepdf.Pdf.new()
        result.pages.extend(pdf.pages[number - 1] for number in selected)
        result.save(output)
    return validate_output(output, len(selected))


def repair_pdf(source: Path, output: Path) -> tuple[int, dict[str, Any]]:
    try:
        with pikepdf.open(source, attempt_recovery=True) as pdf:
            count = len(pdf.pages)
            warnings = list(pdf.get_warnings())
            pdf.save(
                output,
                compress_streams=True,
                object_stream_mode=pikepdf.ObjectStreamMode.generate,
            )
    except pikepdf.PdfError as exc:
        raise LocalPDFError("PDF_CORRUPT", "The PDF could not be repaired.") from exc
    validate_output(output, count)
    return count, {"recovery_warnings": warnings}


def images_to_pdf(inputs: list[Path], output: Path) -> int:
    if not inputs:
        raise LocalPDFError("INPUT_INVALID", "Select at least one image.")
    document = pymupdf.open()
    try:
        for source in inputs:
            try:
                with Image.open(source) as opened:
                    image = ImageOps.exif_transpose(opened).convert("RGB")
                    buffer = io.BytesIO()
                    image.save(buffer, format="JPEG", quality=92)
                    width, height = image.size
            except (OSError, ValueError) as exc:
                raise LocalPDFError("UNSUPPORTED_MEDIA_TYPE", "An image could not be opened.") from exc
            max_width, max_height = A4
            scale = min(max_width / width, max_height / height)
            page_width, page_height = width * scale, height * scale
            page = document.new_page(width=page_width, height=page_height)
            page.insert_image(page.rect, stream=buffer.getvalue())
        document.save(output, garbage=4, deflate=True)
    finally:
        document.close()
    return validate_output(output, len(inputs))


def office_to_pdf(source: Path, output: Path) -> int:
    extension = source.suffix.lower()
    if extension == ".docx":
        _docx_to_pdf(source, output)
    elif extension == ".xlsx":
        _xlsx_to_pdf(source, output)
    elif extension == ".pptx":
        _pptx_to_pdf(source, output)
    elif extension in {".html", ".htm"}:
        _html_to_pdf(source, output)
    else:
        raise LocalPDFError("UNSUPPORTED_MEDIA_TYPE", "This file cannot be converted to PDF.")
    return validate_output(output)


def _font_name() -> str:
    name = "LocalPDFUnicode"
    if name in pdfmetrics.getRegisteredFontNames():
        return name
    candidates = [
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "arial.ttf",
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            pdfmetrics.registerFont(TTFont(name, str(candidate)))
            return name
    return "Helvetica"


def _styles() -> tuple[Any, Any]:
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "LocalPDFBody",
        parent=styles["BodyText"],
        fontName=_font_name(),
        fontSize=10,
        leading=14,
        spaceAfter=5,
    )
    heading = ParagraphStyle(
        "LocalPDFHeading",
        parent=styles["Heading1"],
        fontName=_font_name(),
        fontSize=18,
        leading=22,
        spaceAfter=10,
    )
    return body, heading


def _docx_to_pdf(source: Path, output: Path) -> None:
    document = WordDocument(str(source))
    body, heading = _styles()
    story: list[Any] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            style = heading if paragraph.style and paragraph.style.name.startswith("Heading") else body
            story.append(Paragraph(html.escape(text), style))
    for source_table in document.tables:
        rows = [[Paragraph(html.escape(cell.text), body) for cell in row.cells] for row in source_table.rows]
        if rows:
            table = Table(rows, repeatRows=1)
            table.setStyle(_table_style())
            story.extend([Spacer(1, 4 * mm), table, Spacer(1, 4 * mm)])
    if not story:
        story.append(Paragraph("Empty document", body))
    SimpleDocTemplate(str(output), pagesize=A4, rightMargin=15 * mm, leftMargin=15 * mm).build(story)


def _xlsx_to_pdf(source: Path, output: Path) -> None:
    workbook = load_workbook(source, read_only=True, data_only=True)
    body, heading = _styles()
    story: list[Any] = []
    for sheet_index, sheet in enumerate(workbook.worksheets):
        if sheet_index:
            story.append(PageBreak())
        story.append(Paragraph(html.escape(sheet.title), heading))
        rows: list[list[Any]] = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                rows.append([Paragraph(html.escape(value), body) for value in values])
        if rows:
            available = landscape(A4)[0] - 30 * mm
            widths = [available / max(1, len(rows[0]))] * len(rows[0])
            table = Table(rows, colWidths=widths, repeatRows=1)
            table.setStyle(_table_style())
            story.append(table)
        else:
            story.append(Paragraph("Empty sheet", body))
    workbook.close()
    SimpleDocTemplate(str(output), pagesize=landscape(A4), rightMargin=15 * mm, leftMargin=15 * mm).build(story)


def _pptx_to_pdf(source: Path, output: Path) -> None:
    presentation = Presentation(str(source))
    from reportlab.pdfgen.canvas import Canvas

    width = float(presentation.slide_width or 9_144_000) / 12700
    height = float(presentation.slide_height or 5_143_500) / 12700
    draw = Canvas(str(output), pagesize=(width, height))
    font = _font_name()
    for slide in presentation.slides:
        draw.setFillColor(colors.white)
        draw.rect(0, 0, width, height, stroke=0, fill=1)
        for shape in slide.shapes:
            left = float(shape.left or 0) / 12700
            top = float(shape.top or 0) / 12700
            shape_width = float(shape.width or 0) / 12700
            shape_height = float(shape.height or 0) / 12700
            if getattr(shape, "has_text_frame", False):
                y = height - top - 12
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        size = min(22.0, max(8.0, shape_height / 5))
                        draw.setFillColor(colors.black)
                        draw.setFont(font, size)
                        draw.drawString(left, y, text[:400])
                        y -= size * 1.3
            elif getattr(shape, "shape_type", None) == 13:
                try:
                    draw.drawImage(
                        io.BytesIO(shape.image.blob),
                        left,
                        height - top - shape_height,
                        width=shape_width,
                        height=shape_height,
                        preserveAspectRatio=True,
                    )
                except Exception:
                    continue
        draw.showPage()
    draw.save()


def _html_to_pdf(source: Path, output: Path) -> None:
    soup = BeautifulSoup(source.read_text(encoding="utf-8", errors="replace"), "html.parser")
    body, heading = _styles()
    story: list[Any] = []
    title = soup.title.get_text(" ", strip=True) if soup.title else source.stem
    story.append(Paragraph(html.escape(title), heading))
    for element in soup.find_all(["h1", "h2", "h3", "p", "li", "blockquote", "pre"]):
        text = element.get_text(" ", strip=True)
        if text:
            story.append(Paragraph(html.escape(text), heading if element.name.startswith("h") else body))
    SimpleDocTemplate(str(output), pagesize=A4).build(story)


def _table_style() -> TableStyle:
    return TableStyle(
        [
            ("FONTNAME", (0, 0), (-1, -1), _font_name()),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E6F4F1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ]
    )


def pdf_to_jpg(source: Path, output: Path, dpi: int = 160) -> dict[str, Any]:
    images: list[tuple[str, bytes]] = []
    with pymupdf.open(source) as document:
        for number, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            images.append((f"page-{number:04d}.jpg", pixmap.tobytes("jpeg", jpg_quality=90)))
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, data in images:
            archive.writestr(name, data)
    return {"page_count": len(images), "format": "JPG"}


def pdf_to_word(source: Path, output: Path) -> dict[str, Any]:
    target = WordDocument()
    with pymupdf.open(source) as document:
        for number, page in enumerate(document, start=1):
            target.add_heading(f"Page {number}", level=1)
            text = page.get_text("text").strip()
            target.add_paragraph(text or "[No extractable text]")
            if number != len(document):
                target.add_page_break()  # type: ignore[no-untyped-call]
        count = len(document)
    target.save(str(output))
    return {"page_count": count, "format": "DOCX"}


def pdf_to_powerpoint(source: Path, output: Path, dpi: int = 130) -> dict[str, Any]:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    with pymupdf.open(source) as document:
        for page in document:
            slide = presentation.slides.add_slide(blank)
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            stream = io.BytesIO(pixmap.tobytes("png"))
            slide.shapes.add_picture(stream, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        count = len(document)
    presentation.save(str(output))
    return {"page_count": count, "format": "PPTX"}


def pdf_to_excel(source: Path, output: Path) -> dict[str, Any]:
    workbook = Workbook()
    workbook.remove(workbook.active)
    with pymupdf.open(source) as document:
        for number, page in enumerate(document, start=1):
            sheet = workbook.create_sheet(f"Page {number}")
            text = page.get_text("text")
            for row, line in enumerate(text.splitlines(), start=1):
                cells = [cell.strip() for cell in re.split(r"\t|\s{2,}", line) if cell.strip()]
                for column, value in enumerate(cells or [line], start=1):
                    sheet.cell(row=row, column=column, value=value)
        count = len(document)
    workbook.save(output)
    return {"page_count": count, "format": "XLSX"}


def pdf_to_pdfa(source: Path, output: Path) -> tuple[int, dict[str, Any]]:
    with pikepdf.open(source) as pdf:
        with pdf.open_metadata(set_pikepdf_as_editor=False) as metadata:
            metadata["pdfaid:part"] = "2"
            metadata["pdfaid:conformance"] = "B"
            metadata["dc:format"] = "application/pdf"
        count = len(pdf.pages)
        pdf.save(output, linearize=True, fix_metadata_version=True)
    validate_output(output, count)
    return count, {"profile": "PDF/A-2b", "validation": "metadata-profile"}


def add_page_numbers(source: Path, output: Path, parameters: dict[str, Any]) -> int:
    font_size = float(parameters.get("font_size", 10))
    prefix = str(parameters.get("prefix", ""))
    with pymupdf.open(source) as document:
        count = len(document)
        for number, page in enumerate(document, start=1):
            label = f"{prefix}{number} / {count}"
            rect = pymupdf.Rect(20, page.rect.height - 34, page.rect.width - 20, page.rect.height - 12)
            page.insert_textbox(rect, label, fontsize=font_size, align=TA_CENTER, color=(0.1, 0.1, 0.1))
        document.save(output, garbage=4, deflate=True)
    return validate_output(output, count)


def crop_pdf(source: Path, output: Path, parameters: dict[str, Any]) -> int:
    left = float(parameters.get("left", 0))
    top = float(parameters.get("top", 0))
    right = float(parameters.get("right", 0))
    bottom = float(parameters.get("bottom", 0))
    with pymupdf.open(source) as document:
        for page in document:
            box = page.cropbox
            target = pymupdf.Rect(box.x0 + left, box.y0 + top, box.x1 - right, box.y1 - bottom)
            if target.width < 36 or target.height < 36:
                raise LocalPDFError("INPUT_INVALID", "Crop margins leave no usable page area.")
            page.set_cropbox(target)
        count = len(document)
        document.save(output, garbage=4, deflate=True)
    return validate_output(output, count)


def edit_pdf(source: Path, output: Path, parameters: dict[str, Any]) -> int:
    text = str(parameters.get("text", "")).strip()
    if not text:
        raise LocalPDFError("INPUT_INVALID", "Enter text to add.")
    page_number = int(parameters.get("page_number", 1))
    x = float(parameters.get("x", 72))
    y = float(parameters.get("y", 72))
    font_size = float(parameters.get("font_size", 14))
    with pymupdf.open(source) as document:
        if page_number < 1 or page_number > len(document):
            raise LocalPDFError("PAGE_OUT_OF_BOUNDS", "The selected page does not exist.")
        page = document[page_number - 1]
        rect = pymupdf.Rect(x, y, min(page.rect.width - 12, x + 360), min(page.rect.height - 12, y + 90))
        page.insert_textbox(rect, text, fontsize=font_size, color=(0, 0, 0))
        count = len(document)
        document.save(output, garbage=4, deflate=True)
    return validate_output(output, count)


def add_pdf_form(source: Path, output: Path, parameters: dict[str, Any]) -> int:
    page_number = int(parameters.get("page_number", 1))
    label = str(parameters.get("name", "field_1")).strip() or "field_1"
    x = float(parameters.get("x", 72))
    y = float(parameters.get("y", 72))
    width = float(parameters.get("width", 240))
    height = float(parameters.get("height", 34))
    with pymupdf.open(source) as document:
        if page_number < 1 or page_number > len(document):
            raise LocalPDFError("PAGE_OUT_OF_BOUNDS", "The selected page does not exist.")
        widget = pymupdf.Widget()
        widget.field_name = label
        widget.field_type = pymupdf.PDF_WIDGET_TYPE_TEXT
        widget.field_value = str(parameters.get("value", ""))
        widget.rect = pymupdf.Rect(x, y, x + width, y + height)
        widget.text_fontsize = 11
        widget.field_flags = 0
        document[page_number - 1].add_widget(widget)
        count = len(document)
        document.save(output, garbage=4, deflate=True)
    return validate_output(output, count)


def unlock_pdf(source: Path, output: Path, password: str) -> int:
    try:
        with pikepdf.open(source, password=password) as pdf:
            count = len(pdf.pages)
            pdf.save(output)
    except pikepdf.PasswordError as exc:
        raise LocalPDFError("PASSWORD_INVALID", "The PDF password is incorrect.") from exc
    return validate_output(output, count)


def protect_pdf(source: Path, output: Path, password: str) -> dict[str, Any]:
    if len(password) < 4:
        raise LocalPDFError("INPUT_INVALID", "Use a password with at least 4 characters.")
    with pikepdf.open(source) as pdf:
        count = len(pdf.pages)
        pdf.save(
            output,
            encryption=pikepdf.Encryption(
                owner=password,
                user=password,
                R=6,
                allow=pikepdf.Permissions(extract=False, modify_annotation=False, modify_other=False),
            ),
        )
    if not output.exists() or output.stat().st_size == 0:
        raise LocalPDFError("OUTPUT_VALIDATION_FAILED", "The protected PDF was not created.")
    return {"page_count": count, "encrypted": True}


def compare_pdfs(first: Path, second: Path, output: Path) -> tuple[int, dict[str, Any]]:
    before = _extract_text(first).splitlines()
    after = _extract_text(second).splitlines()
    diff = list(difflib.unified_diff(before, after, fromfile=first.name, tofile=second.name, lineterm=""))
    body, heading = _styles()
    story: list[Any] = [Paragraph("PDF comparison report", heading)]
    story.append(Paragraph(f"Changed lines: {sum(1 for line in diff if line.startswith(('+', '-')) and not line.startswith(('+++', '---')))}", body))
    for line in diff[:2000]:
        story.append(Paragraph(html.escape(line) or "&nbsp;", body))
    if not diff:
        story.append(Paragraph("No textual differences found.", body))
    SimpleDocTemplate(str(output), pagesize=A4).build(story)
    count = validate_output(output)
    return count, {"different": bool(diff), "diff_line_count": len(diff)}


def summarize_pdf(source: Path, output: Path, max_sentences: int = 8) -> tuple[int, dict[str, Any]]:
    text = _extract_text(source)
    sentences = [value.strip() for value in re.split(r"(?<=[.!?])\s+|\n+", text) if len(value.strip()) > 20]
    words = [word.casefold() for word in re.findall(r"[\wÀ-ɏ]+", text) if len(word) > 3]
    frequencies = Counter(words)
    ranked = sorted(
        enumerate(sentences),
        key=lambda item: sum(frequencies[word.casefold()] for word in re.findall(r"[\wÀ-ɏ]+", item[1])) / max(1, len(item[1])),
        reverse=True,
    )[: max(1, max_sentences)]
    summary = " ".join(sentence for _, sentence in sorted(ranked)) or text[:2000] or "No extractable text."
    _text_to_pdf("AI Summary", summary, output)
    count = validate_output(output)
    return count, {"summary": summary[:4000], "sentence_count": len(ranked), "engine": "offline-extractive"}


def translate_pdf(source: Path, output: Path, source_language: str, target_language: str) -> tuple[int, dict[str, Any]]:
    if source_language == target_language:
        translated = _extract_text(source)
    else:
        model = _model_directory(source_language, target_language)
        processor = sentencepiece.SentencePieceProcessor(
            model_proto=(model / "sentencepiece.model").read_bytes()
        )
        translator = ctranslate2.Translator(str(model / "model"), device="cpu")
        sentences = [value.strip() for value in re.split(r"(?<=[.!?])\s+|\n+", _extract_text(source)) if value.strip()]
        batches = [processor.encode(sentence, out_type=str) for sentence in sentences]
        translated_parts: list[str] = []
        for start in range(0, len(batches), 32):
            results = translator.translate_batch(batches[start : start + 32], beam_size=3)
            translated_parts.extend(processor.decode(result.hypotheses[0]) for result in results)
        translated = "\n\n".join(translated_parts) or "No extractable text."
    _text_to_pdf(f"Translation: {source_language.upper()} → {target_language.upper()}", translated, output)
    count = validate_output(output)
    return count, {"source_language": source_language, "target_language": target_language, "characters": len(translated)}


def pdf_to_markdown(source: Path, output: Path) -> dict[str, Any]:
    lines = ["# LocalPDF export", ""]
    with pymupdf.open(source) as document:
        for number, page in enumerate(document, start=1):
            lines.extend([f"## Page {number}", "", page.get_text("text").strip(), ""])
        count = len(document)
    output.write_text("\n".join(lines), encoding="utf-8")
    return {"page_count": count, "format": "Markdown"}


def ocr_pdf(source: Path, output: Path, parameters: dict[str, Any]) -> tuple[int, dict[str, Any]]:
    engine = RapidOCR()
    recognized = 0
    with pymupdf.open(source) as document:
        for page in document:
            pixmap = page.get_pixmap(dpi=180, alpha=False)
            image = np.array(Image.open(io.BytesIO(pixmap.tobytes("png"))).convert("RGB"))
            result = cast(Any, engine(image))
            if result.boxes is None or result.txts is None:
                continue
            scale_x = page.rect.width / pixmap.width
            scale_y = page.rect.height / pixmap.height
            for box, text in zip(result.boxes, result.txts, strict=True):
                x0 = float(min(point[0] for point in box)) * scale_x
                y0 = float(min(point[1] for point in box)) * scale_y
                x1 = float(max(point[0] for point in box)) * scale_x
                y1 = float(max(point[1] for point in box)) * scale_y
                rect = pymupdf.Rect(x0, y0, max(x0 + 2, x1), max(y0 + 2, y1))
                page.insert_text(
                    (rect.x0, max(rect.y0 + 4, rect.y1)),
                    text,
                    fontsize=max(4, min(14, rect.height * 0.8)),
                    render_mode=3,
                )
                recognized += 1
        count = len(document)
        document.save(output, garbage=4, deflate=True)
    validate_output(output, count)
    return count, {"recognized_blocks": recognized, "engine": "RapidOCR"}


def _extract_text(source: Path) -> str:
    with pymupdf.open(source) as document:
        return "\n\n".join(page.get_text("text") for page in document)


def _text_to_pdf(title: str, text: str, output: Path) -> None:
    body, heading = _styles()
    story: list[Any] = [Paragraph(html.escape(title), heading)]
    for paragraph in re.split(r"\n{2,}", text):
        if paragraph.strip():
            story.append(Paragraph(html.escape(paragraph.strip()).replace("\n", "<br/>"), body))
    SimpleDocTemplate(str(output), pagesize=A4, rightMargin=15 * mm, leftMargin=15 * mm).build(story)


def _model_directory(source_language: str, target_language: str) -> Path:
    if (source_language, target_language) not in {("en", "tr"), ("tr", "en")}:
        raise LocalPDFError("INPUT_INVALID", "Offline translation currently supports English and Turkish.")
    roots: Iterable[Path] = (
        Path(os.environ.get("LOCALPDF_MODEL_DIR", "")),
        Path(sys.executable).resolve().parent / "models",
        Path(__file__).resolve().parents[5] / "tools" / "desktop-runtime" / "build" / "models" / "extracted",
    )
    name = f"{source_language}-{target_language}"
    package_name = f"translate-{source_language}_{target_language}-1_5"
    for root in roots:
        candidate = root / name / package_name
        if (candidate / "model" / "model.bin").exists():
            return candidate
    raise LocalPDFError("TOOL_UNAVAILABLE", "The offline translation model is not installed.", status_code=503)
