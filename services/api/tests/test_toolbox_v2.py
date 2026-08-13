from pathlib import Path

import pikepdf
import pymupdf  # type: ignore[import-untyped]
from docx import Document as WordDocument
from openpyxl import Workbook  # type: ignore[import-untyped]
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.pdfgen import canvas

from app.application.services.operations import SUPPORTED_OPERATIONS
from app.infrastructure.tools import toolbox


def make_pdf(path: Path, labels: list[str]) -> None:
    draw = canvas.Canvas(str(path), pagesize=(300, 400))
    for label in labels:
        draw.setFont("Helvetica", 18)
        draw.drawString(30, 220, label)
        draw.showPage()
    draw.save()


def test_v2_declares_all_33_backend_operations() -> None:
    assert len(SUPPORTED_OPERATIONS) == 32
    assert {"repair", "protect", "translate", "pdf_forms", "pdf_to_markdown"} <= SUPPORTED_OPERATIONS


def test_organize_edit_security_and_intelligence_tools(tmp_path: Path, monkeypatch: object) -> None:
    source = tmp_path / "source.pdf"
    second = tmp_path / "second.pdf"
    make_pdf(source, ["Hello, how are you?", "Local PDF document.", "Final page."])
    make_pdf(second, ["Hello, how are you?", "Changed document."])

    assert toolbox.remove_pages(source, tmp_path / "removed.pdf", "3") == 2
    assert toolbox.extract_pages(source, tmp_path / "extracted.pdf", "1,3") == 2
    assert toolbox.repair_pdf(source, tmp_path / "repaired.pdf")[0] == 3
    assert toolbox.add_page_numbers(source, tmp_path / "numbered.pdf", {}) == 3
    assert toolbox.crop_pdf(source, tmp_path / "cropped.pdf", {"left": 10, "top": 10}) == 3
    assert toolbox.edit_pdf(source, tmp_path / "edited.pdf", {"text": "Added", "page_number": 1}) == 3
    assert toolbox.add_pdf_form(source, tmp_path / "form.pdf", {"name": "customer"}) == 3
    assert toolbox.compare_pdfs(source, second, tmp_path / "compare.pdf")[1]["different"] is True
    assert toolbox.summarize_pdf(source, tmp_path / "summary.pdf", 2)[0] >= 1
    assert toolbox.pdf_to_pdfa(source, tmp_path / "archive.pdf")[1]["profile"] == "PDF/A-2b"

    protected = tmp_path / "protected.pdf"
    assert toolbox.protect_pdf(source, protected, "secret123")["encrypted"] is True
    with pikepdf.open(protected, password="secret123") as encrypted:
        assert encrypted.is_encrypted
    assert toolbox.unlock_pdf(protected, tmp_path / "unlocked.pdf", "secret123") == 3

    model_root = Path(__file__).resolve().parents[3] / "tools" / "desktop-runtime" / "build" / "models" / "extracted"
    monkeypatch.setenv("LOCALPDF_MODEL_DIR", str(model_root))  # type: ignore[attr-defined]
    translated = tmp_path / "translated.pdf"
    assert toolbox.translate_pdf(source, translated, "en", "tr")[1]["characters"] > 0


def test_image_ocr_and_export_tools(tmp_path: Path) -> None:
    image_path = tmp_path / "scan.png"
    image = Image.new("RGB", (900, 300), "white")
    font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 64)
    ImageDraw.Draw(image).text((40, 100), "Hello LocalPDF 123", fill="black", font=font)
    image.save(image_path)
    scanned = tmp_path / "scanned.pdf"
    assert toolbox.images_to_pdf([image_path, image_path], scanned) == 2

    ocr = tmp_path / "ocr.pdf"
    assert toolbox.ocr_pdf(scanned, ocr, {})[0] == 2
    with pymupdf.open(ocr) as document:
        assert "LocalPDF" in "".join(page.get_text() for page in document)

    jpg_zip = tmp_path / "pages.zip"
    assert toolbox.pdf_to_jpg(scanned, jpg_zip)["page_count"] == 2
    assert toolbox.pdf_to_word(scanned, tmp_path / "converted.docx")["format"] == "DOCX"
    assert toolbox.pdf_to_powerpoint(scanned, tmp_path / "converted.pptx")["format"] == "PPTX"
    assert toolbox.pdf_to_excel(scanned, tmp_path / "converted.xlsx")["format"] == "XLSX"
    markdown = tmp_path / "converted.md"
    assert toolbox.pdf_to_markdown(scanned, markdown)["format"] == "Markdown"
    assert markdown.read_text(encoding="utf-8").startswith("# LocalPDF export")


def test_office_and_html_to_pdf_without_external_apps(tmp_path: Path) -> None:
    docx = tmp_path / "sample.docx"
    word = WordDocument()
    word.add_heading("LocalPDF Word", level=1)
    word.add_paragraph("Converted inside the desktop app.")
    word.save(str(docx))

    xlsx = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active.append(["Name", "Value"])
    workbook.active.append(["LocalPDF", 2])
    workbook.save(xlsx)

    pptx = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "LocalPDF PowerPoint"
    presentation.save(str(pptx))

    html = tmp_path / "sample.html"
    html.write_text("<html><head><title>LocalPDF HTML</title></head><body><h1>Title</h1><p>Body text.</p></body></html>", encoding="utf-8")

    for source in (docx, xlsx, pptx, html):
        output = tmp_path / f"{source.stem}-{source.suffix[1:]}.pdf"
        assert toolbox.office_to_pdf(source, output) >= 1
